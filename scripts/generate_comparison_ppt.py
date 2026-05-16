"""
Generate a polished Customer Comparison One-Pager PPT highlighting MSIL DA2.8 advantages.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys, io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

OUTPUT = r"C:\Users\mdevarapaga\Downloads\Customer_Comparison_One_Pager_v2.pptx"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x36, 0x5F)
MED_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0x4A, 0x86, 0xC8)
ACCENT_GREEN = RGBColor(0x2D, 0x8C, 0x48)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_AMBER = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE1)
MSIL_HIGHLIGHT = RGBColor(0xE3, 0xF2, 0xFD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_table_cell_format(cell, text, font_size=10, bold=False, color=DARK_GRAY, fill_color=None, alignment=PP_ALIGN.CENTER):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


# ============================================================
# SLIDE 1 — Title & Context
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide1, WHITE)

# Top banner
add_rect(slide1, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BLUE)
add_textbox(slide1, Inches(0.5), Inches(0.15), Inches(8), Inches(0.5),
            "MSIL DA2.8 — Customer Defect Comparison", 28, True, WHITE)
add_textbox(slide1, Inches(0.5), Inches(0.7), Inches(8), Inches(0.4),
            "Benchmarking against 4 comparable IVI programs | April 2026", 14, False, RGBColor(0xBB, 0xDE, 0xFB))

# Left panel — Context
add_rounded_rect(slide1, Inches(0.4), Inches(1.6), Inches(5.8), Inches(5.5), LIGHT_GRAY)
add_textbox(slide1, Inches(0.7), Inches(1.8), Inches(5.2), Inches(0.4),
            "Comparison Scope", 18, True, DARK_BLUE)

context_items = [
    ("Benchmark Programs", "4 comparable IVI programs — 2 Qualcomm, 2 Exynos-based"),
    ("Platform Type", "All platform-based IVI-only programs"),
    ("MSIL Advantage", "First Android 14 project on QC; others on Android 9/11"),
    ("Development Time", "MSIL has the shortest development cycle among all (< 2 years)"),
    ("SOP Status", "2 programs have achieved SOP; others still in development"),
]
y_pos = 2.4
for title, desc in context_items:
    add_textbox(slide1, Inches(0.9), Inches(y_pos), Inches(5), Inches(0.25),
                f"▸ {title}", 12, True, MED_BLUE)
    add_textbox(slide1, Inches(1.1), Inches(y_pos + 0.28), Inches(4.8), Inches(0.35),
                desc, 11, False, DARK_GRAY)
    y_pos += 0.75

# Right panel — Program comparison table
table_data = [
    ["Project", "Android Ver.", "CPU / RAM"],
    ["MSIL DA2.8", "14 (Latest)", "QC6150, 6 GB"],
    ["P1", "14", "Samsung V820, 16 GB"],
    ["P2", "9", "Intel Apollo Lake, 6 GB"],
    ["P3", "11", "QC6150, 4/6 GB"],
    ["P4", "14", "Samsung V820, 16 GB"],
]

tbl_shape = slide1.shapes.add_table(len(table_data), 3, Inches(6.8), Inches(1.6), Inches(6), Inches(3.2))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(1.7)
tbl.columns[2].width = Inches(2.5)

for r_idx, row_data in enumerate(table_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl.cell(r_idx, c_idx)
        if r_idx == 0:
            add_table_cell_format(cell, val, 11, True, WHITE, DARK_BLUE)
        elif r_idx == 1:
            add_table_cell_format(cell, val, 11, True, DARK_BLUE, MSIL_HIGHLIGHT)
        else:
            add_table_cell_format(cell, val, 10, False, DARK_GRAY, WHITE)

# Key callout boxes on right
callout_y = 5.2
callouts = [
    ("✦ First A14 on QC Platform", ACCENT_GREEN, GREEN_BG),
    ("✦ Shortest Development Cycle", ACCENT_GREEN, GREEN_BG),
    ("✦ 2 Programs Already at SOP", MED_BLUE, MSIL_HIGHLIGHT),
]
for text, text_color, bg_color in callouts:
    box = add_rounded_rect(slide1, Inches(6.8), Inches(callout_y), Inches(6), Inches(0.45), bg_color)
    add_textbox(slide1, Inches(7.0), Inches(callout_y + 0.05), Inches(5.5), Inches(0.35),
                text, 13, True, text_color)
    callout_y += 0.55

# ============================================================
# SLIDE 2 — Detailed Metrics Comparison
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)

# Top banner
add_rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide2, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "Defect Metrics Comparison — MSIL DA2.8 vs Peers", 26, True, WHITE)
add_textbox(slide2, Inches(0.5), Inches(0.6), Inches(10), Inches(0.3),
            "Lower is better for all defect percentages", 12, False, RGBColor(0xBB, 0xDE, 0xFB))

# Main comparison table
main_data = [
    ["Project", "CPU / RAM", "Total\nDefects", "TOP-A\nDefects", "TOP-A\n%", "Stability\nDefects", "Stability\n%", "HMI\nDefects", "HMI\n%", "Top 2 Defect\nDomains"],
    ["MSIL DA2.8", "QC6150, 6 GB", "10,469", "2,034", "18%", "255", "2.3%", "5,815", "55.5%", "HMI & Projection"],
    ["P1", "V820, 16 GB", "3,357", "938", "27%", "196", "5.8%", "1,308", "39.1%", "HMI & Foundation SW"],
    ["P2", "Apollo Lake, 6 GB", "27,357", "13,493", "49%", "851", "3.1%", "7,979", "29.2%", "HMI & Systems"],
    ["P3", "QC6150, 4/6 GB", "18,409", "5,702", "36%", "737", "4.0%", "6,208", "33.7%", "HMI & Projection"],
    ["P4", "V820, 16 GB", "7,025", "4,863", "69%", "533", "7.6%", "1,326", "18.8%", "HMI & Systems"],
]

tbl2_shape = slide2.shapes.add_table(len(main_data), 10, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.0))
tbl2 = tbl2_shape.table
col_widths = [1.4, 1.4, 0.9, 0.9, 0.8, 0.9, 0.9, 0.9, 0.8, 1.8]  # ~12.7 total
for i, w in enumerate(col_widths):
    tbl2.columns[i].width = Inches(w)

for r_idx, row_data in enumerate(main_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl2.cell(r_idx, c_idx)
        if r_idx == 0:
            add_table_cell_format(cell, val, 9, True, WHITE, DARK_BLUE)
        elif r_idx == 1:
            # MSIL row — highlight
            add_table_cell_format(cell, val, 10, True, DARK_BLUE, MSIL_HIGHLIGHT)
        else:
            add_table_cell_format(cell, val, 9, False, DARK_GRAY, WHITE)

# Percentage advantage callouts
add_textbox(slide2, Inches(0.3), Inches(4.5), Inches(12.7), Inches(0.4),
            "MSIL DA2.8 — Key Advantages Over Peer Programs", 18, True, DARK_BLUE)

# Advantage cards
cards = [
    {
        "title": "TOP-A Defects",
        "value": "33–74% Lower",
        "detail": "MSIL 18% vs peers 27–69%\nLowest critical defect ratio among all programs",
        "icon": "▼",
        "color": ACCENT_GREEN,
        "bg": GREEN_BG,
    },
    {
        "title": "Stability Defects",
        "value": "26–70% Lower",
        "detail": "MSIL 2.3% vs peers 3.1–7.6%\nBest stability record across all platforms",
        "icon": "▼",
        "color": ACCENT_GREEN,
        "bg": GREEN_BG,
    },
    {
        "title": "Development Time",
        "value": "Shortest Cycle",
        "detail": "< 2 years on newest Android 14\nFirst A14 on Qualcomm platform",
        "icon": "⚡",
        "color": ACCENT_GREEN,
        "bg": GREEN_BG,
    },
    {
        "title": "HMI Defects",
        "value": "Higher (55.5%)",
        "detail": "Due to parallel development\nExpected to normalize post-stabilization",
        "icon": "▲",
        "color": ACCENT_AMBER,
        "bg": AMBER_BG,
    },
]

card_x = 0.3
for card in cards:
    box = add_rounded_rect(slide2, Inches(card_x), Inches(5.0), Inches(3.0), Inches(2.2), card["bg"])
    # Icon + Title
    add_textbox(slide2, Inches(card_x + 0.15), Inches(5.1), Inches(2.7), Inches(0.3),
                f"{card['icon']}  {card['title']}", 13, True, card["color"])
    # Value
    add_textbox(slide2, Inches(card_x + 0.15), Inches(5.45), Inches(2.7), Inches(0.4),
                card["value"], 22, True, card["color"])
    # Detail
    add_textbox(slide2, Inches(card_x + 0.15), Inches(5.95), Inches(2.7), Inches(0.8),
                card["detail"], 10, False, DARK_GRAY)
    card_x += 3.25

# ============================================================
# SLIDE 3 — Conclusions & Summary
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)

# Top banner
add_rect(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
add_textbox(slide3, Inches(0.5), Inches(0.15), Inches(10), Inches(0.5),
            "Conclusions — MSIL DA2.8 Stands Out", 26, True, WHITE)
add_textbox(slide3, Inches(0.5), Inches(0.6), Inches(10), Inches(0.3),
            "Summary of competitive positioning vs peer IVI programs", 12, False, RGBColor(0xBB, 0xDE, 0xFB))

# Scorecard table
score_data = [
    ["Metric", "MSIL DA2.8", "Peer Range", "MSIL vs Best Peer", "Verdict"],
    ["TOP-A Defect %", "18%", "27% – 69%", "33% lower than nearest (P1: 27%)", "✔ MSIL Leads"],
    ["Stability Defect %", "2.3%", "3.1% – 7.6%", "26% lower than nearest (P2: 3.1%)", "✔ Best in Class"],
    ["HMI Defect %", "55.5%", "18.8% – 39.1%", "Higher — parallel dev phase", "⚠ Expected"],
    ["Android Version", "14 (Latest)", "9 / 11 / 14", "On par or ahead", "✔ Leading Edge"],
    ["Development Time", "< 2 years", "2+ years", "Shortest among all", "✔ Fastest"],
    ["SOP Readiness", "On Track", "2 at SOP, 2 not", "Competitive", "✔ On Track"],
]

tbl3_shape = slide3.shapes.add_table(len(score_data), 5, Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.0))
tbl3 = tbl3_shape.table
score_widths = [2.0, 1.5, 2.0, 4.0, 1.8]
for i, w in enumerate(score_widths):
    tbl3.columns[i].width = Inches(w)

verdict_colors = {
    "✔ MSIL Leads": (ACCENT_GREEN, GREEN_BG),
    "✔ Best in Class": (ACCENT_GREEN, GREEN_BG),
    "⚠ Expected": (ACCENT_AMBER, AMBER_BG),
    "✔ Leading Edge": (ACCENT_GREEN, GREEN_BG),
    "✔ Fastest": (ACCENT_GREEN, GREEN_BG),
    "✔ On Track": (ACCENT_GREEN, GREEN_BG),
}

for r_idx, row_data in enumerate(score_data):
    for c_idx, val in enumerate(row_data):
        cell = tbl3.cell(r_idx, c_idx)
        if r_idx == 0:
            add_table_cell_format(cell, val, 11, True, WHITE, DARK_BLUE)
        elif c_idx == 4:  # Verdict column
            v_color, v_bg = verdict_colors.get(val, (DARK_GRAY, WHITE))
            add_table_cell_format(cell, val, 10, True, v_color, v_bg)
        elif c_idx == 1:  # MSIL column
            add_table_cell_format(cell, val, 10, True, DARK_BLUE, MSIL_HIGHLIGHT)
        else:
            add_table_cell_format(cell, val, 10, False, DARK_GRAY, WHITE)

# Bottom conclusions
add_textbox(slide3, Inches(0.3), Inches(4.6), Inches(12.7), Inches(0.4),
            "Key Takeaways", 18, True, DARK_BLUE)

conclusions = [
    ("✔", "33–74% fewer TOP-A defects than peers — MSIL has the lowest critical defect ratio at 18% vs 27–69%.", ACCENT_GREEN),
    ("✔", "26–70% fewer stability defects — MSIL at 2.3% is the best among all compared programs (next best: 3.1%).", ACCENT_GREEN),
    ("✔", "Shortest development cycle (< 2 years) while being the first program on the newest Android 14 + QC6150 platform.", ACCENT_GREEN),
    ("⚠", "HMI defect % is higher (55.5%) due to parallel development of HMI features — expected to normalize post-stabilization phase.", ACCENT_AMBER),
    ("✔", "Despite being on the latest and most complex platform (A14), MSIL delivers superior defect metrics compared to mature A9/A11 programs.", ACCENT_GREEN),
]

y = 5.1
for icon, text, color in conclusions:
    # Icon
    add_textbox(slide3, Inches(0.5), Inches(y), Inches(0.4), Inches(0.35),
                icon, 14, True, color)
    # Text
    add_textbox(slide3, Inches(0.9), Inches(y), Inches(12), Inches(0.35),
                text, 12, False, DARK_GRAY)
    y += 0.42

# Footer
add_rect(slide3, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), DARK_BLUE)
add_textbox(slide3, Inches(0.5), Inches(7.13), Inches(12), Inches(0.3),
            "MSIL DA2.8 | Confidential | April 2026", 10, False, WHITE, PP_ALIGN.CENTER)

# Save
prs.save(OUTPUT)
print(f"Saved to: {OUTPUT}")
print("3 slides: Title & Context | Detailed Metrics | Conclusions")
